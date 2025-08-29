import io
import os
import tempfile
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from slide_schema import SlidePlan
from llm_router import plan_slides, clean_slide_plan
from typing import Optional

app = FastAPI(title="Auto PPT Generator", version="1.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/health")
def health():
    return {"ok": True}

@app.post("/generate")
async def generate_ppt(
    text: str = Form(...),
    provider: str = Form(...),    # "openai" | "anthropic" | "gemini"
    api_key: str = Form(...),     # not stored, just used for this request
    tone: Optional[str] = Form(None),
    template: Optional[UploadFile] = File(None),  # .pptx or .potx
):
    if not text.strip():
        raise HTTPException(status_code=400, detail="Text is empty")

    # 1) Plan slides with the chosen LLM (never store api_key)
    try:
        plan: SlidePlan = plan_slides(provider, api_key, text, tone)
        # ✅ Clean plan before generating PPT
        plan = clean_slide_plan(plan)
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"LLM error: {e}")

    # 2) Load template (preserves theme/backgrounds/layout pictures)
    if template is not None:
        if not (template.filename.endswith(".pptx") or template.filename.endswith(".potx")):
            raise HTTPException(status_code=400, detail="Template must be .pptx or .potx")
        tmp_tpl = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(template.filename)[1])
        tmp_tpl.write(await template.read())
        tmp_tpl.flush()
        prs = Presentation(tmp_tpl.name)
        tmp_tpl.close()
        os.unlink(tmp_tpl.name)
    else:
        prs = Presentation()  # default theme

    # 3) Choose a common layout: Title + Content if available, else Title Only
    layout_idx_title_content = 1 if len(prs.slide_layouts) > 1 else 0
    layout_idx_title_only = 5 if len(prs.slide_layouts) > 5 else 0

    # 4) Build slides
    for item in plan.slides:
        layout = prs.slide_layouts[layout_idx_title_content]
        slide = prs.slides.add_slide(layout)

        # Set title
        slide.shapes.title.text = item.title

        # Body placeholder can vary by template
        body_placeholder = None
        for shp in slide.placeholders:
            if shp.placeholder_format.type == 2:  # BODY
                body_placeholder = shp
                break
        if body_placeholder is None:
            # fall back to a title-only layout + add textbox
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx_title_only])
            slide.shapes.title.text = item.title
            body = slide.shapes.add_textbox(
                prs.slide_width * 0.08,
                prs.slide_height * 0.25,
                prs.slide_width * 0.84,
                prs.slide_height * 0.6,
            )
            tf = body.text_frame
        else:
            tf = body_placeholder.text_frame
            tf.clear()

        # Bullets
        for i, bullet in enumerate(item.bullets):
            if i == 0 and not tf.paragraphs[0].text:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            # Keep font size respectful of template; set only if missing
            if p.runs and p.runs[0].font.size is None:
                p.runs[0].font.size = Pt(18)

        # Speaker notes
        if item.notes:
            notes = slide.notes_slide if slide.has_notes_slide else slide.notes_slide
            notes.notes_text_frame.text = item.notes

    # 5) Save output
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(out.name)
    out.flush()
    return FileResponse(out.name, filename="generated.pptx")

@app.post("/plan")
async def plan_only(
    text: str = Form(...),
    provider: str = Form(...),
    api_key: str = Form(...),
    tone: Optional[str] = Form(None),
):
    plan = plan_slides(provider, api_key, text, tone)
    # ✅ Clean plan before returning JSON too (optional, keeps preview & PPT consistent)
    plan = clean_slide_plan(plan)
    return JSONResponse(plan.model_dump())
